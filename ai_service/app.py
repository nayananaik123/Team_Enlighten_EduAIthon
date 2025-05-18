# ai_service/app.py
from flask import Flask, request, jsonify, render_template, send_from_directory, redirect
from flask_cors import CORS
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords, wordnet
from nltk.tag import pos_tag
from string import punctuation
import random
import os
import re
import base64
import io
import PyPDF2
from docx import Document
from pptx import Presentation
import spacy
import gensim.downloader as api
from concurrent.futures import ThreadPoolExecutor
import multiprocessing
from googletrans import Translator

# Download required NLTK data
try:
    nltk.download('punkt', quiet=True)
    nltk.download('averaged_perceptron_tagger', quiet=True)
    nltk.download('stopwords', quiet=True)
    nltk.download('wordnet', quiet=True)
    nltk.download('omw-1.4', quiet=True)  # Open Multilingual Wordnet
except Exception as e:
    print(f"Error downloading NLTK data: {str(e)}")

# Initialize NLTK components
try:
    from nltk.corpus import wordnet as wn
    from nltk.corpus import stopwords
    STOPWORDS = set(stopwords.words('english'))
except Exception as e:
    print(f"Error initializing NLTK components: {str(e)}")
    STOPWORDS = set()

# Load spaCy model with increased max_length
try:
    nlp = spacy.load('en_core_web_sm')
    nlp.max_length = 5000000  # Increased from 3M to 5M characters
    
    # Only enable components we actually need
    enabled_pipes = ['tagger', 'attribute_ruler']
    disabled_pipes = [pipe for pipe in nlp.pipe_names if pipe not in enabled_pipes]
    nlp.disable_pipes(*disabled_pipes)
except Exception as e:
    print(f"Error loading spaCy model: {str(e)}")
    raise

# Set number of threads for parallel processing
NUM_THREADS = max(2, multiprocessing.cpu_count() - 1)  # Ensure at least 2 threads

# Initialize word vectors cache with a size limit
word_vectors_cache = {}
MAX_CACHE_SIZE = 10000  # Limit cache size to prevent memory issues

# Initialize word2vec model
try:
    word2vec_model = api.load('word2vec-google-news-300')
except Exception as e:
    print(f"Warning: Could not load word2vec model. Using WordNet only for similarity. Error: {str(e)}")
    word2vec_model = None

app = Flask(__name__, static_folder='../frontend', template_folder='../frontend')
CORS(app, resources={r"/*": {"origins": "*"}})

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/static/<path:path>')
def send_static(path):
    return send_from_directory('static', path)

@app.route('/profile', methods=['GET'])
def profile_page():
    # Dummy user info for now
    return render_template('profile.html', username='DemoUser', email='demo@example.com')

def extract_text_from_pdf(pdf_content):
    try:
        pdf_file = io.BytesIO(pdf_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        print(f"Error extracting PDF text: {str(e)}")
        return None

def extract_text_from_docx(docx_content):
    try:
        doc = Document(io.BytesIO(docx_content))
        text = []
        for paragraph in doc.paragraphs:
            text.append(paragraph.text)
        return '\n'.join(text)
    except Exception as e:
        print(f"Error extracting DOCX text: {str(e)}")
        return None

def extract_text_from_pptx(pptx_content):
    try:
        prs = Presentation(io.BytesIO(pptx_content))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    except Exception as e:
        print(f"Error extracting PPTX text: {str(e)}")
        return None

def clean_text(text):
    if not text:
        return None
        
    # Remove special characters and normalize whitespace
    text = re.sub(r'[^\w\s.,!?-]', ' ', text)
    text = re.sub(r'\s+', ' ', text)
    text = text.strip()
    
    # If text is too short or seems like garbage, return error
    if len(text.split()) < 10 or not any(c.isalpha() for c in text):
        return None
        
    return text

def process_large_text(text, chunk_size=900000):
    """Process large text by breaking it into chunks."""
    chunks = [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]
    processed_chunks = []
    
    # Process chunks in parallel using multiple CPU cores
    with ThreadPoolExecutor(max_workers=NUM_THREADS) as executor:
        processed_chunks = list(executor.map(lambda chunk: nlp(chunk), chunks))
    
    return processed_chunks

def get_important_sentences(text, num_sentences=10):
    print(f"Processing text of length {len(text)} for important sentences") # Debug log
    
    # For very large texts, first do a quick split to get manageable chunks
    sentences = sent_tokenize(text)
    print(f"Found {len(sentences)} total sentences") # Debug log
    
    # Filter out very short or problematic sentences
    sentences = [s.strip() for s in sentences if len(s.strip()) > 20 and len(s.strip()) < 300]
    print(f"After filtering, have {len(sentences)} valid sentences") # Debug log
    
    # For very large texts, take a strategic sample from different sections
    if len(sentences) > 500:
        # Calculate section sizes - divide into 10 sections for better coverage
        num_sections = 10
        section_size = len(sentences) // num_sections
        
        # Take samples from each section to ensure coverage across the document
        samples = []
        for i in range(num_sections):
            start_idx = i * section_size
            end_idx = start_idx + section_size
            section = sentences[start_idx:end_idx]
            
            # Take more samples from each section for large texts
            sample_size = max(20, len(section) // 10)  # Increased minimum samples per section
            if section:  # Check if section is not empty
                section_samples = random.sample(section, min(sample_size, len(section)))
                samples.extend(section_samples)
        
        sentences = samples
        print(f"After sampling, selected {len(sentences)} sentences for processing") # Debug log
    
    # Score sentences
    sentence_scores = {}
    for sentence in sentences:
        score = 0
        words = word_tokenize(sentence.lower())
        
        # Calculate sentence score based on multiple factors
        # 1. Length score - prefer medium length sentences
        length_score = 1.0
        if 15 <= len(words) <= 35:
            length_score = 1.5
        
        # 2. Content score - check for important indicators
        content_score = 1.0
        important_phrases = ['key', 'important', 'main', 'significant', 'essential', 'critical', 'fundamental']
        if any(phrase in sentence.lower() for phrase in important_phrases):
            content_score = 1.5
        
        # 3. Structure score - well-formed sentences
        structure_score = 1.0
        if sentence[0].isupper() and sentence[-1] in '.!?':
            structure_score = 1.2
        
        # 4. Word complexity score
        complex_words = sum(1 for word in words if len(word) > 6)
        complexity_score = 1.0 + (complex_words / len(words))
        
        # Combine scores
        final_score = length_score * content_score * structure_score * complexity_score
        sentence_scores[sentence] = final_score
    
    # Get top sentences ensuring diversity
    sorted_sentences = sorted(sentence_scores.items(), key=lambda x: x[1], reverse=True)
    selected_sentences = []
    seen_content = set()
    
    for sentence, _ in sorted_sentences:
        # Create a content signature to avoid similar sentences
        content_words = set(word.lower() for word in word_tokenize(sentence) 
                          if word.lower() not in stopwords.words('english'))
        content_sig = ' '.join(sorted(content_words))
        
        # Check if we already have a similar sentence
        if not any(len(content_words.intersection(set(seen.split()))) / len(content_words) > 0.6 
                  for seen in seen_content):
            selected_sentences.append(sentence)
            seen_content.add(content_sig)
            if len(selected_sentences) >= num_sentences:
                break
    
    print(f"Final selection: {len(selected_sentences)} sentences") # Debug log
    return selected_sentences

def get_important_words(text, num_words=20):
    try:
        print("Starting key concepts extraction with text length:", len(text))
        
        # Process text with spaCy
        doc = nlp(text[:min(len(text), nlp.max_length)])
        
        # Track word frequencies with POS and custom weights
        word_freq = {}
        
        # Process each token
        for token in doc:
            # Only consider content words
            if (not token.is_stop and 
                not token.is_punct and 
                token.is_alpha and 
                len(token.text) > 3):
                
                word = token.text.lower()
                
                # Weight by POS tag
                weight = 1.0
                if token.pos_ in ['NOUN', 'PROPN']:
                    weight = 3.0
                elif token.pos_ == 'VERB':
                    weight = 2.0
                elif token.pos_ == 'ADJ':
                    weight = 1.5
                
                # Additional weights for technical terms
                if any(c.isupper() for c in token.text):  # Acronyms
                    weight *= 1.5
                if len(token.text) > 6:  # Longer words
                    weight *= 1.2
                
                # Update frequency
                word_freq[word] = word_freq.get(word, 0) + weight
        
        # Sort by frequency
        sorted_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)
        
        # Get top words, avoiding duplicates
        unique_words = []
        seen_roots = set()
        
        for word, _ in sorted_words:
            # Simple stemming
            root = word[:4]
            if root not in seen_roots:
                unique_words.append(word)
                seen_roots.add(root)
            if len(unique_words) >= num_words:
                break
        
        print(f"Extracted {len(unique_words)} key concepts:", unique_words)
        return unique_words

    except Exception as e:
        print(f"Error in key concepts extraction: {str(e)}")
        return ["Error extracting key concepts"]

def generate_flashcards(text):
    try:
        print("Starting flashcard generation with text length:", len(text))
        flashcards = []
        
        # Get important sentences and words
        important_sentences = get_important_sentences(text, num_sentences=15)
        important_words = get_important_words(text, num_words=15)
        
        print(f"Found {len(important_words)} important words")
        
        # Create definition-based flashcards
        for word in important_words:
            # Find context using string operations
            sentences = sent_tokenize(text)
            word_context = ""
            
            # Find a sentence containing the word
            for sent in sentences:
                if word.lower() in sent.lower():
                    word_context = sent
                    break
            
            # Get word definition
            synsets = wn.synsets(word)
            if synsets:
                word_definition = synsets[0].definition()
                
                flashcard = {
                    'term': word,
                    'definition': word_definition,
                    'context': word_context,
                    'type': 'definition'
                }
                flashcards.append(flashcard)
                print(f"Created flashcard for word: {word}")
        
        # Create concept-based flashcards
        for sentence in important_sentences[:5]:
            words = word_tokenize(sentence)
            tagged = pos_tag(words)
            
            # Find subject and verb
            subject = None
            verb = None
            
            for word, tag in tagged:
                if not subject and tag.startswith('NN'):
                    subject = word
                elif not verb and tag.startswith('VB'):
                    verb = word
                if subject and verb:
                    break
            
            if subject and verb:
                question = f"What does {subject} {verb}?"
                flashcard = {
                    'term': question,
                    'definition': sentence,
                    'type': 'concept'
                }
                flashcards.append(flashcard)
                print(f"Created concept flashcard: {question}")
        
        print(f"Generated total {len(flashcards)} flashcards")
        return flashcards
        
    except Exception as e:
        print(f"Error in flashcard generation: {str(e)}")
        return []

def generate_quiz_questions(text):
    try:
        print("Starting quiz generation")
        questions = []
        
        # Process text into sentences
        sentences = sent_tokenize(text)
        print(f"Found {len(sentences)} sentences")
        
        # Filter and select sentences
        valid_sentences = [s for s in sentences if len(s.split()) > 5 and len(s) < 200]
        selected_sentences = random.sample(valid_sentences, min(len(valid_sentences), 50))
        
        # Process selected sentences
        processed_sentences = []
        for sentence in selected_sentences:
            try:
                doc = nlp(sentence)
                processed_sentences.append((sentence, doc))
            except Exception as e:
                print(f"Error processing sentence: {str(e)}")
                continue
        
        # Generate questions
        for sentence, doc in processed_sentences:
            try:
                # Extract key terms
                key_terms = [(token.text, token.pos_) for token in doc 
                            if not token.is_stop and not token.is_punct and len(token.text) > 3]
                
                if not key_terms:
                    continue
                
                # Select a random term to create a fill-in-the-blank question
                term, pos = random.choice(key_terms)
                
                # Create the question
                question = f"Complete this statement: {sentence.replace(term, '_____')}"
                
                # Generate options
                options = [term]  # Correct answer
                
                # Add distractors
                distractors = []
                for other_sent, other_doc in processed_sentences:
                    if other_sent != sentence:
                        similar_terms = [t.text for t in other_doc if t.pos_ == pos and t.text != term]
                        distractors.extend(similar_terms)
                
                if distractors:
                    options.extend(random.sample(distractors, min(3, len(distractors))))
                
                # If we need more options
                while len(options) < 4:
                    options.append(random.choice(['other', 'none', 'unknown', 'N/A']))
                
                # Shuffle options
                random.shuffle(options)
                
                # Find the index of the correct answer
                correctAnswer = options.index(term)
                
                questions.append({
                    "question": question,
                    "options": options,
                    "correctAnswer": correctAnswer
                })
                
                if len(questions) >= 5:  # Limit to 5 questions
                    break
            
            except Exception as e:
                print(f"Error generating question: {str(e)}")
                continue
        
        print(f"Generated {len(questions)} quiz questions")
        print("Generated quiz questions:", questions)  # Debug log
        return questions
    
    except Exception as e:
        print(f"Error in quiz generation: {str(e)}")
        return []

def get_key_points(sentence):
    # Convert statement to question
    words = word_tokenize(sentence)
    tagged = pos_tag(words)
    
    # Find main subject and verb
    subject = None
    verb = None
    for word, tag in tagged:
        if not subject and tag.startswith('NN'):
            subject = word
        elif not verb and tag.startswith('VB'):
            verb = word
        if subject and verb:
            break
    
    if subject and verb:
        return f"What {verb} {subject}?"
    else:
        return f"What is the main point about {words[0]}?"

def get_similar_words(word):
    similar_words = set()
    
    # Try word2vec model first
    if word2vec_model:
        try:
            similar = word2vec_model.most_similar(word.lower(), topn=3)
            similar_words.update(w for w, _ in similar)
        except Exception as e:
            print(f"Word2vec error for word '{word}': {str(e)}")
    
    # Use WordNet as backup
    try:
        synsets = wn.synsets(word)
        for syn in synsets:
            for lemma in syn.lemmas():
                if lemma.name() != word:
                    similar_words.add(lemma.name())
                if len(similar_words) >= 3:
                    break
            if len(similar_words) >= 3:
                break
    except Exception as e:
        print(f"WordNet error for word '{word}': {str(e)}")
    
    # If no similar words found, return some generic alternatives
    if not similar_words:
        domain_words = {
            'NOUN': ["system", "process", "data", "application", "software"],
            'VERB': ["process", "analyze", "compute", "generate", "implement"],
            'ADJ': ["efficient", "automated", "integrated", "secure", "reliable"]
        }
        similar_words = set(random.sample(domain_words['NOUN'], min(3, len(domain_words['NOUN']))))
    
    return list(similar_words)[:3]

def modify_sentence_for_false_statement(sentence):
    words = word_tokenize(sentence)
    tagged = pos_tag(words)
    
    # Find words to potentially replace
    candidates = []
    for i, (word, tag) in enumerate(tagged):
        if tag.startswith(('NN', 'VB', 'JJ')) and len(word) > 3:
            candidates.append((i, word))
    
    if not candidates:
        return None
    
    # Choose a random word to replace
    idx, word = random.choice(candidates)
    similar_words = get_similar_words(word)
    
    if not similar_words:
        return None
    
    # Replace the word with a similar one
    replacement = random.choice(similar_words)
    words[idx] = replacement
    
    return ' '.join(words)

@app.route('/process', methods=['POST'])
def process_text():
    try:
        if 'file' not in request.files:
            data = request.get_json()
            if not data or 'file' not in data:
                return jsonify({"error": "No file provided"}), 400
                
            # Handle base64 encoded file
            file_data = data['file']
            content = base64.b64decode(file_data['content'])
            file_type = file_data['type'].lower()
        else:
            file = request.files['file']
            content = file.read()
            file_type = file.filename.split('.')[-1].lower()
        
        # Extract text based on file type
        if file_type == 'pdf':
            text = extract_text_from_pdf(content)
        elif file_type in ['doc', 'docx']:
            text = extract_text_from_docx(content)
        elif file_type in ['ppt', 'pptx']:
            text = extract_text_from_pptx(content)
        elif file_type == 'txt':
            text = content.decode('utf-8')
        else:
            return jsonify({"error": "Unsupported file type"}), 400
        
        # Clean and validate text
        text = clean_text(text)
        if not text:
            return jsonify({"error": "Could not extract valid text from file"}), 400
        
        print("Processing text of length:", len(text)) # Debug log
        
        # Process text in parallel
        with ThreadPoolExecutor(max_workers=NUM_THREADS) as executor:
            # Start all tasks
            print("Starting key concepts extraction...") # Debug log
            keywords_future = executor.submit(get_important_words, text)
            
            print("Starting flashcard generation...") # Debug log
            flashcards_future = executor.submit(generate_flashcards, text)
            
            print("Starting quiz generation...") # Debug log
            quiz_future = executor.submit(generate_quiz_questions, text)
            
            # Get results
            keywords = keywords_future.result()
            print(f"Generated {len(keywords)} key concepts:", keywords) # Debug log
            
            flashcards = flashcards_future.result()
            print(f"Generated {len(flashcards)} flashcards") # Debug log
            
            quiz_questions = quiz_future.result()
            print(f"Generated {len(quiz_questions)} quiz questions") # Debug log
            
            # Format quiz questions to ensure correct types
            formatted_quiz_questions = []
            for q in quiz_questions:
                if 'question' in q and 'options' in q and 'correctAnswer' in q:
                    formatted_quiz_questions.append({
                        'question': str(q['question']),
                        'options': [str(opt) for opt in q['options']],
                        'correctAnswer': int(q['correctAnswer'])
                    })
            print("Formatted quiz questions:", formatted_quiz_questions) # Debug log
            
            # Create summary from important sentences
            important_sentences = get_important_sentences(text, num_sentences=3)
            summary = ' '.join(important_sentences)
        
        # Format flashcards for frontend
        formatted_flashcards = []
        for card in flashcards:
            if card['type'] == 'definition':
                formatted_flashcards.append({
                    "front": f"What is the meaning of '{card['term']}'?",
                    "back": f"Definition: {card['definition']}\n\nContext: {card['context']}"
                })
            else:
                formatted_flashcards.append({
                    "front": card['term'],
                    "back": card['definition']
                })
        
        # Format key concepts as a list of strings
        formatted_keywords = [str(keyword) for keyword in keywords]
        
        response_data = {
            "text": text,
            "summary": summary,
            "keywords": formatted_keywords,  # Use formatted keywords
            "flashcards": formatted_flashcards,
            "quizQuestions": formatted_quiz_questions
        }
        
        print("Sending response with keywords:", response_data["keywords"]) # Debug log
        print("Sending response with flashcards:", response_data["flashcards"]) # Debug log
        
        return jsonify(response_data)
        
    except Exception as e:
        print(f"Error processing file: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/translate', methods=['POST'])
def api_translate():
    data = request.json
    text = data.get('text')
    source = data.get('source')
    target = data.get('target')
    if not text or not target:
        return {'error': 'Missing text or target language'}, 400
    translator = Translator()
    result = translator.translate(text, src=source, dest=target)
    return {'translatedText': result.text}

if __name__ == '__main__':
    app.run(port=5001, debug=True)